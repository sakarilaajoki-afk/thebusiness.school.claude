"""
Spec checker for thebusiness.school teaching resources.

Why this exists. Every resource states exam-board facts in its own prose, and
there was no single place recording what is true. So a correct code in one file
sat next to a wrong one in another, and nothing broke when someone wrote a code
that does not exist. The errors were only ever found by a human re-reading, and
a human reading 160 files misses things.

This turns the recurring findings into a check that runs in seconds.

    python _spec-check/check.py

Exit code 0 = clean, 1 = something to look at. Add --quiet for CI use.

It reads _spec-check/board-facts.json, which records what was verified, from
which board document, on what date. If the checker and the registry disagree
with a resource, open the board document in the registry before changing
anything: the registry can go stale too.
"""
import json, os, re, sys, subprocess

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
RES = os.path.join(ROOT, "resources")
FACTS = json.load(open(os.path.join(HERE, "board-facts.json"), encoding="utf-8"))
Q = FACTS["qualifications"]

SKIP_DIRS = ("_backup", "_varmuuskopio", "_VARMUUSKOPIO", "_ennen", "_uncorrected", "_spec-check")

findings = []


def add(sev, path, msg, detail=""):
    findings.append((sev, os.path.relpath(path, ROOT), msg, detail))


# ---------------------------------------------------------------- text loading
def text_of(path):
    ext = path.lower().rsplit(".", 1)[-1]
    try:
        if ext in ("html", "htm", "txt", "md"):
            s = open(path, encoding="utf-8", errors="replace").read()
            s = re.sub(r"<script.*?</script>|<style.*?</style>", " ", s, flags=re.S)
            return re.sub(r"<[^>]+>", " ", s)
        if ext == "pdf":
            import fitz
            d = fitz.open(path)
            return "\n".join(d[i].get_text() for i in range(d.page_count))
        if ext == "docx":
            import docx
            d = docx.Document(path)
            return "\n".join(p.text for p in d.paragraphs)
        if ext == "pptx":
            from pptx import Presentation
            pr = Presentation(path)
            out = []
            for sl in pr.slides:
                for sh in sl.shapes:
                    if sh.has_text_frame:
                        out.append(sh.text_frame.text)
            return "\n".join(out)
    except Exception as e:
        add("SKIP", path, "could not open", repr(e)[:90])
    return ""


def files():
    for dp, dn, fn in os.walk(RES):
        dn[:] = [d for d in dn if not any(s in d for s in SKIP_DIRS)]
        for f in fn:
            if f.lower().rsplit(".", 1)[-1] in ("html", "pdf", "docx", "pptx"):
                yield os.path.join(dp, f)


# ---------------------------------------------------------------- the checks
def check_aqa_gcse_codes(path, t):
    """AQA 8132 numbers content to three levels only. 3.6.2.1 does not exist."""
    if "8132" not in t:
        return
    bad = sorted(set(re.findall(r"\b3\.\d+\.\d+\.\d+\b", t)))
    for c in bad:
        add("ERROR", path, "AQA GCSE 8132 code has four levels and cannot exist", c)
    valid = set(Q["aqa-gcse-business-8132"]["valid_codes"])
    for c in sorted(set(re.findall(r"(?:8132|AQA GCSE)[^.\n]{0,40}?(3\.\d(?:\.\d)?)", t))):
        if c not in valid:
            add("ERROR", path, "AQA GCSE 8132 code not in the specification", c)


def check_aqa_alevel_codes(path, t):
    """AQA 7132 section 3.1 ends at 3.1.3.

    Attribution matters and getting it wrong produces noise. The same string
    '3.1.4' is a VALID code for AQA GCSE 8132 and a VALID code for AQA 7138
    (Financial management). It is only wrong when it is attributed to 7132.
    So only look at codes that follow a 7132 label and stop at the next board
    label, rather than anywhere in a file that happens to mention 7132.
    """
    if "7132" not in t:
        return
    for seg in re.finditer(
            r"7132\b((?:[^.]|\.\d){0,80})", t):
        chunk = seg.group(1)
        # stop at the next board label so we do not swallow another board's codes
        chunk = re.split(r"8132|9BS0|1BS0|H4\d\d|J204|0450|0455|7138|Edexcel|OCR|Cambridge",
                         chunk)[0]
        for c in sorted(set(re.findall(r"\b3\.1\.([4-9])\b", chunk))):
            add("ERROR", path, "AQA 7132 section 3.1 ends at 3.1.3, so this code does not exist",
                "3.1." + c + "  (markets and market research are 3.3.2, positioning 3.3.3, added value 3.4.1)")


def check_currency(path, t):
    """Outgoing codes must carry a note. A September 2026 Year 12 is not on them."""
    note = re.search(r"7138|H436|outgoing|summer 2027|last exams|September 2026", t, re.I)
    for code, key in (("7132", "aqa-alevel-business-7132"), ("H431", "ocr-alevel-business-h431")):
        if code in t and Q[key].get("currency_note_required") and not note:
            add("WARN", path, "names an outgoing specification with no currency note", code)
    if "0450" in t and not re.search(r"0264|withdraw|final (?:exam|series)|November 2026", t, re.I):
        add("WARN", path, "names Cambridge 0450 with no note that it becomes 0264", "0450")
    if "0455" in t and not re.search(r"2027|2028|2029", t):
        add("WARN", path, "names Cambridge 0455 without saying which syllabus version", "0455")


def check_excluded_content(path, t):
    """Content the board says it will not assess, taught as if it will be."""
    low = t.lower()
    if "8132" in t:
        if re.search(r"calculate the break[- ]even", low):
            add("ERROR", path, "AQA 8132 excludes this",
                'AQA: "Students will not be expected to draw break-even charts or use the break-even formula."')
        if "contribution" in low and "8132" in t:
            add("WARN", path, "'contribution' does not appear in the AQA 8132 specification", "contribution")
    if "0450" in t:
        for ex in Q["cambridge-igcse-business-0450"]["excluded_content"]:
            if "PED" in ex["quote"] and re.search(r"calculat\w+ (?:the )?(?:price elasticity|PED)", low):
                add("ERROR", path, "Cambridge 0450 excludes this", ex["quote"])


def check_ao_attribution(path, t):
    """Pearson publishes no per-question AO split for the 9BS0 20-marker."""
    if "9BS0" not in t:
        return
    if re.search(r"20[ -]mark", t) and re.search(r"AO split|Knowledge \d+, Application", t, re.I):
        if not re.search(r"\bours\b|our own|author's own|inference|not published", t, re.I):
            add("WARN", path, "attributes a per-question AO split to Pearson for the 9BS0 20-mark question",
                Q["edexcel-alevel-business-9bs0"]["question_forms"]["20-mark"]["ao_split_note"])


def check_btec_tech_award(path, t):
    """The 2022 Tech Award is marked in Mark Bands, not Pass/Merit/Distinction criteria."""
    if "Tech Award" not in t or "Enterprise" not in t:
        return
    if not re.search(r"2022", t):
        return
    if re.search(r"[A-C]\.[12][PMD]\d", t):
        add("ERROR", path, "criteria codes like A.2P1 do not exist in the 2022 Tech Award",
            "the 2022 internal components use Mark Bands 0 to 4 against Learning outcomes")
    if re.search(r"Pass,? Merit and Distinction criteria", t, re.I):
        add("ERROR", path, "the 2022 Tech Award has no Pass/Merit/Distinction criteria",
            Q["btec-tech-award-enterprise-2022"]["internal_assessment_model"]["warning"])
    if re.search(r"learning aims?", t, re.I) and "Tech Award" in t:
        add("WARN", path, "the 2022 Tech Award uses 'learning outcomes', not 'learning aims'", "learning aim")


def check_ib(path, t):
    """The IB toolkit is not Unit 6."""
    if "Business Management" in t and re.search(r"Toolkit \(6\.\d\)|Unit 6", t):
        add("ERROR", path, "the IB Business Management guide has five units and no Unit 6",
            Q["ib-business-management"]["toolkit"])


def check_ocr_levels(path, t):
    """OCR uses four levels on its 20-mark Paper 3 questions, not three everywhere."""
    if "H431" not in t:
        return
    if re.search(r"three levels", t, re.I) and not re.search(r"four", t, re.I):
        add("WARN", path, "OCR H431 uses three levels on the 15-marker but four on the 20-mark Paper 3 questions",
            "Evaluate and Discuss both reach 20 marks on H431/03")


CHECKS = [check_aqa_gcse_codes, check_aqa_alevel_codes, check_currency,
          check_excluded_content, check_ao_attribution, check_btec_tech_award,
          check_ib, check_ocr_levels]


def main():
    quiet = "--quiet" in sys.argv
    n = 0
    for p in files():
        t = text_of(p)
        if not t:
            continue
        n += 1
        t = re.sub(r"[ \t]+", " ", t)
        for c in CHECKS:
            c(p, t)

    order = {"ERROR": 0, "WARN": 1, "SKIP": 2}
    findings.sort(key=lambda f: (order.get(f[0], 3), f[1]))
    errs = sum(1 for f in findings if f[0] == "ERROR")
    warns = sum(1 for f in findings if f[0] == "WARN")

    if not quiet or findings:
        print("spec-check: %d files read, %d errors, %d warnings\n" % (n, errs, warns))
        for sev, path, msg, detail in findings:
            print("%-5s %s" % (sev, path))
            print("      %s" % msg)
            if detail:
                print("      %s" % detail[:150])
        if not findings:
            print("nothing to look at.")
    print("\nRegistry last verified %s. Re-open the board documents in "
          "board-facts.json before trusting any of this." % FACTS["_last_verified"])
    return 1 if errs else 0


if __name__ == "__main__":
    sys.exit(main())
